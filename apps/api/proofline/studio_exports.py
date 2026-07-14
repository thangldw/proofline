from __future__ import annotations

import csv
import hashlib
import io
import json
import re
import textwrap
import zipfile
from dataclasses import dataclass
from html import escape

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt

from .models import SourceVersion, StudioArtifact, StudioCitation


class StudioExportError(ValueError):
    pass


@dataclass(frozen=True)
class StudioExport:
    filename: str
    media_type: str
    content: bytes


def _safe_filename(value: str) -> str:
    value = re.sub(r"[^A-Za-z0-9._-]+", "-", value).strip("-.").lower()
    return value[:80] or "proofline-studio"


def _ordered_citations(artifact: StudioArtifact) -> list[StudioCitation]:
    return sorted(artifact.citations, key=lambda citation: citation.ordinal)


def verify_export_provenance(artifact: StudioArtifact, version: SourceVersion) -> None:
    if version.id != artifact.source_version_id or version.source_id != artifact.source_id:
        raise StudioExportError("artifact_source_version_mismatch")
    for citation in _ordered_citations(artifact):
        if citation.source_version_id != version.id or citation.source_id != version.source_id:
            raise StudioExportError("citation_source_version_mismatch")
        if not 0 <= citation.start_offset < citation.end_offset <= len(version.content):
            raise StudioExportError("citation_offsets_invalid")
        exact = version.content[citation.start_offset : citation.end_offset]
        if exact != citation.quote:
            raise StudioExportError("citation_quote_mismatch")
        if hashlib.sha256(exact.encode("utf-8")).hexdigest() != citation.quote_hash:
            raise StudioExportError("citation_hash_mismatch")


def _manifest(artifact: StudioArtifact, version: SourceVersion, source_title: str) -> dict:
    return {
        "schema": "proofline.studio-export.v1",
        "artifact": {
            "id": artifact.id,
            "kind": artifact.kind,
            "title": artifact.title,
            "generation_method": artifact.generation_method,
            "content": artifact.content_json,
        },
        "immutable_source": {
            "id": artifact.source_id,
            "title": source_title,
            "version_id": version.id,
            "version_number": version.version_number,
            "content_sha256": version.content_hash,
        },
        "citations": [
            {
                "ordinal": citation.ordinal,
                "source_id": citation.source_id,
                "source_version_id": citation.source_version_id,
                "start_offset": citation.start_offset,
                "end_offset": citation.end_offset,
                "start_line": citation.start_line,
                "end_line": citation.end_line,
                "quote": citation.quote,
                "quote_sha256": citation.quote_hash,
            }
            for citation in _ordered_citations(artifact)
        ],
    }


def _markdown(artifact: StudioArtifact) -> bytes:
    lines = [f"# {artifact.title}", "", artifact.content_json.get("summary", ""), ""]
    for item in artifact.content_json.get("items", []):
        lines.extend(
            [
                f"## {item['title']}",
                "",
                item.get("answer", item.get("body", "")),
                "",
                f"Evidence: citation {item['citation'] + 1}",
                "",
            ]
        )
    return "\n".join(lines).encode("utf-8")


def _csv_table(artifact: StudioArtifact) -> bytes:
    output = io.StringIO(newline="")
    writer = csv.writer(output)
    columns = artifact.content_json.get("columns") or ["Title", "Body", "Citation"]
    writer.writerow(columns)
    for item in artifact.content_json.get("items", []):
        fallback = [item["title"], item.get("body", ""), item["citation"] + 1]
        writer.writerow(item.get("cells") or fallback)
    return output.getvalue().encode("utf-8-sig")


def _presentation(artifact: StudioArtifact) -> bytes:
    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    title_slide = deck.slides.add_slide(deck.slide_layouts[0])
    title_slide.shapes.title.text = artifact.title
    title_slide.placeholders[1].text = artifact.content_json.get("summary", "")
    for item in artifact.content_json.get("items", []):
        slide = deck.slides.add_slide(deck.slide_layouts[1])
        slide.shapes.title.text = item["title"]
        frame = slide.placeholders[1].text_frame
        frame.text = item.get("body", "")
        evidence = frame.add_paragraph()
        evidence.text = f"Evidence {item['citation'] + 1} — see evidence-manifest.json"
        evidence.font.size = Pt(14)
    output = io.BytesIO()
    deck.save(output)
    return output.getvalue()


def _infographic(artifact: StudioArtifact) -> bytes:
    image = Image.new("RGB", (1600, 1000), "#f7f9ff")
    draw = ImageDraw.Draw(image)
    font = ImageFont.load_default(size=28)
    title_font = ImageFont.load_default(size=46)
    draw.rounded_rectangle(
        (50, 45, 1550, 955), radius=36, fill="#ffffff", outline="#cbd8ff", width=3
    )
    draw.text((100, 90), artifact.title, fill="#15213b", font=title_font)
    colors = ("#e8efff", "#e1f5e8", "#fff0ed", "#f3e8f4", "#f7f3df", "#def5fa")
    for index, item in enumerate(artifact.content_json.get("items", [])[:6]):
        column = index % 2
        row = index // 2
        x = 100 + column * 740
        y = 190 + row * 235
        draw.rounded_rectangle((x, y, x + 680, y + 190), radius=22, fill=colors[index])
        text = f"{index + 1}. {item['title']}\n" + "\n".join(
            textwrap.wrap(item.get("body", ""), width=48)[:3]
        )
        draw.multiline_text((x + 30, y + 25), text, fill="#20304f", font=font, spacing=8)
    output = io.BytesIO()
    image.save(output, format="PNG", optimize=True)
    return output.getvalue()


def _storyboard_html(artifact: StudioArtifact) -> bytes:
    sections = "".join(
        f"<section><h2>{escape(item['title'])}</h2><p>{escape(item.get('body', ''))}</p>"
        f"<small>Evidence {item['citation'] + 1}</small></section>"
        for item in artifact.content_json.get("items", [])
    )
    html = (
        "<!doctype html><meta charset='utf-8'><title>"
        + escape(artifact.title)
        + "</title><style>body{font:18px system-ui;max-width:960px;margin:40px auto;"
        "background:#f7f9ff}"
        "section{background:white;padding:32px;margin:24px;border-radius:24px}small{color:#315fc8}</style>"
        + sections
    )
    return html.encode("utf-8")


def build_studio_export(
    artifact: StudioArtifact, version: SourceVersion, source_title: str
) -> StudioExport:
    verify_export_provenance(artifact, version)
    stem = _safe_filename(artifact.title)
    files: dict[str, bytes] = {
        "evidence-manifest.json": json.dumps(
            _manifest(artifact, version, source_title), ensure_ascii=False, indent=2, sort_keys=True
        ).encode("utf-8"),
        f"{stem}.md": _markdown(artifact),
    }
    if artifact.kind == "presentation":
        files[f"{stem}.pptx"] = _presentation(artifact)
    elif artifact.kind == "infographic":
        files[f"{stem}.png"] = _infographic(artifact)
    elif artifact.kind == "data_table":
        files[f"{stem}.csv"] = _csv_table(artifact)
    elif artifact.kind == "video_overview":
        files[f"{stem}.html"] = _storyboard_html(artifact)
    elif artifact.kind == "audio_overview":
        files[f"{stem}-narration.txt"] = _markdown(artifact)

    archive = io.BytesIO()
    with zipfile.ZipFile(archive, "w", compression=zipfile.ZIP_DEFLATED) as bundle:
        for name, content in sorted(files.items()):
            bundle.writestr(name, content)
    return StudioExport(
        filename=f"{stem}-evidence-package.zip",
        media_type="application/zip",
        content=archive.getvalue(),
    )
